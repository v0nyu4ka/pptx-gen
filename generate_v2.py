#!/usr/bin/env python3
"""PPTX Generator v2 — light theme, better layout, images work"""

import argparse, json, os, sys, requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HF_BASE = "https://router.huggingface.co"
TEXT_MODEL = "meta-llama/Llama-3.1-8B-Instruct"
IMAGE_MODEL = "black-forest-labs/FLUX.1-schnell"

THEMES = {
    "light": {
        "bg": RGBColor(0xff, 0xff, 0xff),
        "title": RGBColor(0x2d, 0x2d, 0x2d),
        "bullet": RGBColor(0x44, 0x44, 0x44),
        "accent": RGBColor(0x43, 0x7e, 0xc4),
        "subtitle": RGBColor(0x66, 0x66, 0x66),
    },
    "warm": {
        "bg": RGBColor(0xff, 0xf8, 0xf0),
        "title": RGBColor(0x3d, 0x2b, 0x1f),
        "bullet": RGBColor(0x55, 0x44, 0x33),
        "accent": RGBColor(0xe8, 0x8d, 0x43),
        "subtitle": RGBColor(0x88, 0x77, 0x66),
    },
    "dark": {
        "bg": RGBColor(0x1a, 0x1a, 0x2e),
        "title": RGBColor(0xff, 0xff, 0xff),
        "bullet": RGBColor(0xcc, 0xcc, 0xcc),
        "accent": RGBColor(0x66, 0x7e, 0xea),
        "subtitle": RGBColor(0x99, 0x99, 0xbb),
    }
}


def hf_text(api_key, prompt, max_tokens=1024):
    resp = requests.post(f"{HF_BASE}/v1/chat/completions", headers={
        "Authorization": f"Bearer {api_key}", "Content-Type": "application/json"
    }, json={"model": TEXT_MODEL, "messages": [{"role": "user", "content": prompt}],
             "max_tokens": max_tokens, "temperature": 0.7})
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]


def hf_image(api_key, prompt):
    resp = requests.post(f"{HF_BASE}/hf-inference/models/{IMAGE_MODEL}", headers={
        "Authorization": f"Bearer {api_key}", "Content-Type": "application/json"
    }, json={"inputs": prompt})
    resp.raise_for_status()
    return resp.content


def generate_slides(api_key, topic, num_slides=6):
    prompt = f"""Generate a presentation about: "{topic}"

Create exactly {num_slides} slides. Return ONLY valid JSON array.
Each slide:
- "title": short slide title
- "bullets": 3-4 concise bullet points
- "image_prompt": vivid English image prompt, bright and colorful, professional photo style. NO text in image.
- "emoji": one relevant emoji for the slide

First slide = title slide (title + subtitle in bullets).
Last slide = summary/conclusion.

Return ONLY JSON array."""
    raw = hf_text(api_key, prompt, 2048)
    start, end = raw.find('['), raw.rfind(']') + 1
    if start == -1 or end == 0:
        print(f"Parse error:\n{raw}", file=sys.stderr); sys.exit(1)
    return json.loads(raw[start:end])


def create_pptx(slides_data, images, output_path, theme_name="warm"):
    theme = THEMES.get(theme_name, THEMES["warm"])
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = theme["bg"]

        # Image on right side
        if i < len(images) and images[i] and i > 0:
            try:
                img_stream = BytesIO(images[i])
                pic = slide.shapes.add_picture(img_stream, Inches(9), Inches(0.5), Inches(6.5), Inches(8))
                # Round corners effect - just position nicely
            except Exception as e:
                print(f"  img error slide {i+1}: {e}", file=sys.stderr)

        # Title slide
        if i == 0:
            # Center image as background if available
            if images and images[0]:
                try:
                    img_stream = BytesIO(images[0])
                    slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(16), Inches(9))
                    # Add overlay
                    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(9))
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = theme["bg"]
                    from pptx.oxml.ns import qn
                    # Make semi-transparent
                except:
                    pass

            emoji = sd.get("emoji", "✨")
            ebox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(1.5))
            p = ebox.text_frame.paragraphs[0]
            p.text = emoji
            p.font.size = Pt(72)
            p.alignment = PP_ALIGN.CENTER

            tbox = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(12), Inches(2))
            tf = tbox.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd["title"]
            p.font.size = Pt(48); p.font.color.rgb = theme["title"]; p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            bullets = sd.get("bullets", [])
            if bullets:
                sbox = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(10), Inches(2))
                tf = sbox.text_frame; tf.word_wrap = True
                for j, b in enumerate(bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(22); p.font.color.rgb = theme["subtitle"]
                    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(8)
            continue

        # Content slides
        emoji = sd.get("emoji", "📌")

        # Emoji + Title
        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(1.2))
        tf = tbox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{emoji}  {sd['title']}"
        p.font.size = Pt(36); p.font.color.rgb = theme["title"]; p.font.bold = True

        # Accent line
        line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.7), Inches(4), Inches(0.06))
        line.fill.solid(); line.fill.fore_color.rgb = theme["accent"]
        line.line.fill.background()

        # Bullets
        bullets = sd.get("bullets", [])
        bbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7.8), Inches(6))
        tf = bbox.text_frame; tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"●  {b}"
            p.font.size = Pt(22); p.font.color.rgb = theme["bullet"]
            p.space_after = Pt(20)

        # Slide number
        nbox = slide.shapes.add_textbox(Inches(14.5), Inches(8.2), Inches(1), Inches(0.5))
        p = nbox.text_frame.paragraphs[0]
        p.text = str(i); p.font.size = Pt(14); p.font.color.rgb = theme["subtitle"]
        p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("topic")
    parser.add_argument("-o", "--output", default="presentation.pptx")
    parser.add_argument("-n", "--slides", type=int, default=6)
    parser.add_argument("-k", "--api-key", default=os.environ.get("HF_API_KEY"))
    parser.add_argument("-t", "--theme", choices=["light", "warm", "dark"], default="warm")
    parser.add_argument("--no-images", action="store_true")
    args = parser.parse_args()
    if not args.api_key:
        print("Need HF API key (-k or HF_API_KEY)", file=sys.stderr); sys.exit(1)

    print(f"📝 Generating: {args.topic}")
    slides = generate_slides(args.api_key, args.topic, args.slides)
    print(f"✅ {len(slides)} slides")

    images = []
    if not args.no_images:
        for i, s in enumerate(slides):
            ip = s.get("image_prompt", "")
            if ip:
                print(f"🎨 Image {i+1}/{len(slides)}: {ip[:50]}...")
                try: images.append(hf_image(args.api_key, ip))
                except Exception as e:
                    print(f"  ⚠️ {e}", file=sys.stderr); images.append(None)
            else: images.append(None)

    print("📦 Building...")
    create_pptx(slides, images, args.output, args.theme)
    print(f"✅ {args.output}")

if __name__ == "__main__":
    main()
