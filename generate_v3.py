#!/usr/bin/env python3
"""PPTX Generator v3 — polished layout, proper image placement"""

import argparse, json, os, sys, requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HF_BASE = "https://router.huggingface.co"
TEXT_MODEL = "meta-llama/Llama-3.1-8B-Instruct"
IMAGE_MODEL = "black-forest-labs/FLUX.1-schnell"

THEMES = {
    "light": {
        "bg": RGBColor(0xf8, 0xf9, 0xfa),
        "title": RGBColor(0x1a, 0x1a, 0x2e),
        "bullet": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x43, 0x7e, 0xc4),
        "subtitle": RGBColor(0x66, 0x66, 0x66),
        "card_bg": RGBColor(0xff, 0xff, 0xff),
        "num_color": RGBColor(0xdd, 0xdd, 0xdd),
    },
    "warm": {
        "bg": RGBColor(0xfe, 0xf9, 0xf4),
        "title": RGBColor(0x2c, 0x1e, 0x15),
        "bullet": RGBColor(0x44, 0x33, 0x22),
        "accent": RGBColor(0xe0, 0x7a, 0x2f),
        "subtitle": RGBColor(0x88, 0x77, 0x66),
        "card_bg": RGBColor(0xff, 0xff, 0xff),
        "num_color": RGBColor(0xee, 0xdd, 0xcc),
    },
    "dark": {
        "bg": RGBColor(0x0f, 0x0f, 0x1a),
        "title": RGBColor(0xf0, 0xf0, 0xf5),
        "bullet": RGBColor(0xbb, 0xbb, 0xcc),
        "accent": RGBColor(0x7c, 0x6a, 0xf7),
        "subtitle": RGBColor(0x88, 0x88, 0xaa),
        "card_bg": RGBColor(0x1a, 0x1a, 0x2e),
        "num_color": RGBColor(0x33, 0x33, 0x44),
    }
}


def hf_text(api_key, prompt, max_tokens=1024):
    resp = requests.post(f"{HF_BASE}/v1/chat/completions", headers={
        "Authorization": f"Bearer {api_key}", "Content-Type": "application/json"
    }, json={"model": TEXT_MODEL, "messages": [{"role": "user", "content": prompt}],
             "max_tokens": max_tokens, "temperature": 0.7})
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]


def hf_image(api_key, prompt):
    resp = requests.post(f"{HF_BASE}/hf-inference/models/{IMAGE_MODEL}", headers={
        "Authorization": f"Bearer {api_key}", "Content-Type": "application/json"
    }, json={"inputs": prompt})
    resp.raise_for_status()
    return resp.content


def add_rounded_rect(slide, left, top, width, height, fill_rgb, radius=Inches(0.15)):
    """Add a rounded rectangle shape"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def generate_slides(api_key, topic, num_slides=6):
    prompt = f"""Generate a presentation about: "{topic}"

Create exactly {num_slides} slides. Return ONLY valid JSON array.
Each slide:
- "title": short punchy slide title (max 6 words)
- "bullets": array of exactly 3 bullet points. Each bullet max 12 words, informative and concise
- "image_prompt": vivid English prompt for AI image. Must be: bright, colorful, professional photography, no text/words in image, 4k quality
- "emoji": one relevant emoji

Slide 1 = title slide. Slide {num_slides} = conclusion/summary.
Return ONLY the JSON array, no markdown."""
    raw = hf_text(api_key, prompt, 2048)
    start, end = raw.find('['), raw.rfind(']') + 1
    if start == -1 or end == 0:
        print(f"Parse error:\n{raw}", file=sys.stderr); sys.exit(1)
    return json.loads(raw[start:end])


def create_pptx(slides_data, images, output_path, theme_name="warm"):
    t = THEMES.get(theme_name, THEMES["warm"])
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for i, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = t["bg"]
        emoji = sd.get("emoji", "✨")

        if i == 0:
            # ===== TITLE SLIDE =====
            # Large image as background
            if images and images[0]:
                try:
                    img_stream = BytesIO(images[0])
                    slide.shapes.add_picture(img_stream, Inches(0), Inches(0), Inches(16), Inches(9))
                    # Dark overlay for readability
                    ov = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(9))
                    ov.fill.solid()
                    ov.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    ov.line.fill.background()
                    # Set opacity via XML
                    pass
                except Exception as e:
                    print(f"  title img: {e}", file=sys.stderr)

            # Emoji
            ebox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(14), Inches(1.2))
            p = ebox.text_frame.paragraphs[0]
            p.text = emoji; p.font.size = Pt(64); p.alignment = PP_ALIGN.CENTER

            # Title
            tbox = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(13), Inches(2.5))
            tf = tbox.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd["title"]
            p.font.size = Pt(52); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

            # Subtitle bullets
            bullets = sd.get("bullets", [])
            if bullets:
                sbox = slide.shapes.add_textbox(Inches(3), Inches(5.8), Inches(10), Inches(2))
                tf = sbox.text_frame; tf.word_wrap = True
                for j, b in enumerate(bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = b; p.font.size = Pt(20); p.alignment = PP_ALIGN.CENTER
                    p.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc); p.space_after = Pt(6)
            continue

        # ===== CONTENT SLIDES =====
        
        # Big slide number in background
        nbg = slide.shapes.add_textbox(Inches(12.5), Inches(0.5), Inches(3), Inches(3))
        p = nbg.text_frame.paragraphs[0]
        p.text = f"{i:02d}"; p.font.size = Pt(120); p.font.color.rgb = t["num_color"]
        p.font.bold = True; p.alignment = PP_ALIGN.RIGHT

        # Image card on right
        if i < len(images) and images[i]:
            try:
                img_stream = BytesIO(images[i])
                img = Image.open(img_stream)
                img_stream.seek(0)
                # Card background
                add_rounded_rect(slide, Inches(9.2), Inches(1), Inches(6.2), Inches(7), t["card_bg"])
                # Image
                pic = slide.shapes.add_picture(img_stream, Inches(9.4), Inches(1.2), Inches(5.8), Inches(5.8))
            except Exception as e:
                print(f"  img slide {i+1}: {e}", file=sys.stderr)

        # Accent bar on left
        bar = slide.shapes.add_shape(1, Inches(0.6), Inches(0.6), Inches(0.08), Inches(7.8))
        bar.fill.solid(); bar.fill.fore_color.rgb = t["accent"]; bar.line.fill.background()

        # Emoji + Title
        tbox = slide.shapes.add_textbox(Inches(1.2), Inches(1), Inches(7.5), Inches(1.5))
        tf = tbox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{emoji}  {sd['title']}"
        p.font.size = Pt(36); p.font.color.rgb = t["title"]; p.font.bold = True

        # Bullets in cards
        bullets = sd.get("bullets", [])
        for j, b in enumerate(bullets):
            y = Inches(2.8 + j * 1.7)
            # Bullet card
            card = add_rounded_rect(slide, Inches(1.2), y, Inches(7.5), Inches(1.4), t["card_bg"])
            
            # Bullet number circle
            cbox = slide.shapes.add_textbox(Inches(1.5), y + Inches(0.25), Inches(0.8), Inches(0.8))
            p = cbox.text_frame.paragraphs[0]
            p.text = str(j + 1); p.font.size = Pt(24); p.font.color.rgb = t["accent"]
            p.font.bold = True; p.alignment = PP_ALIGN.CENTER

            # Bullet text
            bbox = slide.shapes.add_textbox(Inches(2.5), y + Inches(0.2), Inches(5.8), Inches(1))
            tf = bbox.text_frame; tf.word_wrap = True
            tf.paragraphs[0].text = b
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.color.rgb = t["bullet"]

        # Page number bottom right
        pbox = slide.shapes.add_textbox(Inches(14.5), Inches(8.3), Inches(1), Inches(0.4))
        p = pbox.text_frame.paragraphs[0]
        p.text = f"{i}/{len(slides_data)-1}"; p.font.size = Pt(12)
        p.font.color.rgb = t["subtitle"]; p.alignment = PP_ALIGN.RIGHT

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("topic")
    parser.add_argument("-o", "--output", default="presentation.pptx")
    parser.add_argument("-n", "--slides", type=int, default=6)
    parser.add_argument("-k", "--api-key", default=os.environ.get("HF_API_KEY"))
    parser.add_argument("-t", "--theme", choices=["light", "warm", "dark"], default="warm")
    parser.add_argument("--no-images", action="store_true")
    args = parser.parse_args()
    if not args.api_key:
        print("Need HF API key", file=sys.stderr); sys.exit(1)

    print(f"📝 {args.topic}")
    slides = generate_slides(args.api_key, args.topic, args.slides)
    print(f"✅ {len(slides)} slides")

    images = []
    if not args.no_images:
        for i, s in enumerate(slides):
            ip = s.get("image_prompt", "")
            if ip:
                print(f"🎨 {i+1}/{len(slides)}: {ip[:60]}...")
                try: images.append(hf_image(args.api_key, ip))
                except Exception as e: print(f"  ⚠️ {e}", file=sys.stderr); images.append(None)
            else: images.append(None)

    print("📦 Building...")
    create_pptx(slides, images, args.output, args.theme)
    print(f"✅ {args.output}")

if __name__ == "__main__":
    main()
