#!/usr/bin/env python3
"""
PPTX Generator using HuggingFace APIs
- Text generation for slide content
- Image generation for slide backgrounds/illustrations
"""

import argparse
import json
import os
import sys
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HF_BASE = "https://router.huggingface.co"
TEXT_MODEL = "meta-llama/Llama-3.1-8B-Instruct"
IMAGE_MODEL = "black-forest-labs/FLUX.1-schnell"


def hf_text(api_key, prompt, max_tokens=1024):
    """Generate text using HF chat completions API"""
    url = f"{HF_BASE}/v1/chat/completions"
    resp = requests.post(url, headers={
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }, json={
        "model": TEXT_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": max_tokens,
        "temperature": 0.7
    })
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]


def hf_image(api_key, prompt):
    """Generate image using HF inference API"""
    url = f"{HF_BASE}/hf-inference/models/{IMAGE_MODEL}"
    resp = requests.post(url, headers={
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }, json={"inputs": prompt})
    resp.raise_for_status()
    return resp.content


def generate_slide_content(api_key, topic, num_slides=6):
    """Ask LLM to generate structured slide content"""
    prompt = f"""Generate a presentation about: "{topic}"

Create exactly {num_slides} slides. Return ONLY valid JSON array, no other text.
Each slide object must have:
- "title": slide title (short)
- "bullets": array of 3-4 bullet points (concise, informative)  
- "image_prompt": English prompt for AI image generation that would fit this slide (descriptive, visual)

First slide should be a title slide with the presentation title and subtitle in bullets.
Last slide should be a summary/conclusion.

Return ONLY the JSON array, nothing else."""

    raw = hf_text(api_key, prompt, max_tokens=2048)
    
    # Extract JSON from response
    start = raw.find('[')
    end = raw.rfind(']') + 1
    if start == -1 or end == 0:
        print(f"Failed to parse LLM response:\n{raw}", file=sys.stderr)
        sys.exit(1)
    
    return json.loads(raw[start:end])


def create_pptx(slides_data, images, output_path):
    """Build .pptx from structured data + images"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Color scheme
    bg_color = RGBColor(0x1a, 0x1a, 0x2e)
    title_color = RGBColor(0xff, 0xff, 0xff)
    bullet_color = RGBColor(0xcc, 0xcc, 0xcc)
    accent_color = RGBColor(0x66, 0x7e, 0xea)
    
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        
        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add image on right side if available
        if i < len(images) and images[i]:
            try:
                img_stream = BytesIO(images[i])
                img = Image.open(img_stream)
                img_stream.seek(0)
                
                # Place image on right half, slightly transparent overlay
                slide.shapes.add_picture(
                    img_stream,
                    left=Inches(8.5),
                    top=Inches(0),
                    width=Inches(7.5),
                    height=Inches(9)
                )
                
                # Add semi-transparent overlay on image side for readability
                overlay = slide.shapes.add_shape(
                    1,  # rectangle
                    Inches(8.5), Inches(0),
                    Inches(7.5), Inches(9)
                )
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
                overlay.line.fill.background()
            except Exception as e:
                print(f"Warning: failed to add image to slide {i+1}: {e}", file=sys.stderr)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.8) if i > 0 else Inches(2.5),
            Inches(14), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(44) if i == 0 else Pt(36)
        p.font.color.rgb = title_color
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER
        
        # Accent line under title
        if i > 0:
            line = slide.shapes.add_shape(
                1,  # rectangle
                Inches(0.8), Inches(2.2),
                Inches(3), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = accent_color
            line.line.fill.background()
        
        # Bullets
        bullets = slide_data.get("bullets", [])
        if bullets:
            body_top = Inches(2.6) if i > 0 else Inches(4.5)
            body_box = slide.shapes.add_textbox(
                Inches(0.8), body_top,
                Inches(7.5), Inches(5.5)
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            
            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(22) if i > 0 else Pt(24)
                p.font.color.rgb = bullet_color if i > 0 else accent_color
                p.space_after = Pt(16)
                p.alignment = PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER
                if i > 0:
                    p.level = 0
        
        # Slide number
        if i > 0:
            num_box = slide.shapes.add_textbox(
                Inches(14.5), Inches(8.2),
                Inches(1), Inches(0.5)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i)
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p.alignment = PP_ALIGN.RIGHT
    
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate PPTX presentations using HuggingFace AI")
    parser.add_argument("topic", help="Presentation topic")
    parser.add_argument("-o", "--output", default="presentation.pptx", help="Output file path")
    parser.add_argument("-n", "--slides", type=int, default=6, help="Number of slides")
    parser.add_argument("-k", "--api-key", default=os.environ.get("HF_API_KEY"), help="HuggingFace API key")
    parser.add_argument("--no-images", action="store_true", help="Skip image generation")
    args = parser.parse_args()
    
    if not args.api_key:
        print("Error: provide HF API key via -k or HF_API_KEY env var", file=sys.stderr)
        sys.exit(1)
    
    print(f"📝 Generating content for: {args.topic}")
    slides_data = generate_slide_content(args.api_key, args.topic, args.slides)
    print(f"✅ Generated {len(slides_data)} slides")
    
    images = []
    if not args.no_images:
        for i, slide in enumerate(slides_data):
            img_prompt = slide.get("image_prompt", "")
            if img_prompt:
                print(f"🎨 Generating image {i+1}/{len(slides_data)}: {img_prompt[:50]}...")
                try:
                    img_data = hf_image(args.api_key, img_prompt)
                    images.append(img_data)
                except Exception as e:
                    print(f"  ⚠️ Failed: {e}", file=sys.stderr)
                    images.append(None)
            else:
                images.append(None)
    
    print(f"📦 Building PPTX...")
    output = create_pptx(slides_data, images, args.output)
    print(f"✅ Saved to: {output}")


if __name__ == "__main__":
    main()
